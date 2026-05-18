# ╔══════════════════════════════════════════════════════════════╗
# ║           DocuMind AI  —  Intelligent Document Q&A           ║
# ║   Multi-doc · PDF/DOCX/PPTX/MD/TXT · Citations · Web fallback║
# ╚══════════════════════════════════════════════════════════════╝
#
# pip install -r requirements.txt
# Set GOOGLE_API_KEY environment variable before running.

import os
from pathlib import Path
from typing import List

import gradio as gr
from langchain_community.document_loaders import (
    TextLoader,
    PyPDFLoader,
    Docx2txtLoader,
    UnstructuredPowerPointLoader,
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.schema import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_community.tools import DuckDuckGoSearchRun

# ═══════════════════════════════════════════════════════════════
# 1 · GLOBALS  (LLM + embeddings kept exactly as original)
# ═══════════════════════════════════════════════════════════════
embeddings   = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
llm          = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.7)
web_search   = DuckDuckGoSearchRun()
vectorstore  = None
doc_registry = []

SPLITTER = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

FILE_ICONS = {
    ".pdf":  "PDF", ".docx": "DOC", ".doc":  "DOC",
    ".pptx": "PPT", ".ppt":  "PPT", ".md":   "MD",  ".txt":  "TXT",
}

# ═══════════════════════════════════════════════════════════════
# 2 · DOCUMENT LOADING
# ═══════════════════════════════════════════════════════════════
def _pptx_fallback(path: str) -> List[Document]:
    from pptx import Presentation
    prs, docs = Presentation(path), []
    for i, slide in enumerate(prs.slides):
        txt = "\n".join(
            s.text for s in slide.shapes if hasattr(s, "text") and s.text.strip()
        )
        if txt.strip():
            docs.append(Document(page_content=txt,
                                 metadata={"source": path, "page": i + 1}))
    return docs


def load_file(path: str) -> List[Document]:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return PyPDFLoader(path).load()
    elif ext in (".docx", ".doc"):
        return Docx2txtLoader(path).load()
    elif ext in (".pptx", ".ppt"):
        try:
            return UnstructuredPowerPointLoader(path).load()
        except Exception:
            return _pptx_fallback(path)
    elif ext in (".md", ".txt"):
        return TextLoader(path, encoding="utf-8").load()
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def process_files(files):
    global vectorstore, doc_registry
    if not files:
        return _html_doc_list([]), "⚠  No files selected."

    all_chunks, doc_registry = [], []
    for f in files:
        fname = Path(f.name).name
        ext   = Path(f.name).suffix.lower()
        try:
            docs   = load_file(f.name)
            chunks = SPLITTER.split_documents(docs)
            for i, c in enumerate(chunks):
                c.metadata.update({"file_name": fname,
                                   "file_type": ext, "chunk_id": i})
            all_chunks.extend(chunks)
            doc_registry.append({"name": fname, "type": ext,
                                  "chunks": len(chunks), "status": "ok"})
        except Exception as exc:
            doc_registry.append({"name": fname, "type": ext,
                                  "chunks": 0, "status": f"err: {exc}"})

    if all_chunks:
        vectorstore = FAISS.from_documents(all_chunks, embeddings)

    ok           = sum(1 for d in doc_registry if d["status"] == "ok")
    total_chunks = sum(d["chunks"] for d in doc_registry)
    msg = (f"✓  {ok} of {len(files)} file(s) ready"
           f"  ·  {total_chunks} text chunks indexed")
    return _html_doc_list(doc_registry), msg


# ═══════════════════════════════════════════════════════════════
# 3 · HTML BUILDERS
# ═══════════════════════════════════════════════════════════════
def _badge(ext: str) -> str:
    label = FILE_ICONS.get(ext, "FILE")
    colors = {
        "PDF": "#dc2626", "DOC": "#2563eb",
        "PPT": "#ea580c", "MD":  "#16a34a", "TXT": "#7c3aed",
    }
    bg = colors.get(label, "#64748b")
    return (f'<span style="background:{bg};color:#fff;font-size:.62rem;'
            f'font-weight:700;padding:2px 8px;border-radius:4px;'
            f'letter-spacing:.07em">{label}</span>')


def _html_doc_list(docs) -> str:
    if not docs:
        return '<p class="dim">No documents loaded yet.</p>'
    rows = ""
    for d in docs:
        ok    = d["status"] == "ok"
        icon  = "✓" if ok else "✗"
        color = "#16a34a" if ok else "#dc2626"
        rows += f"""
        <div class="doc-row">
          {_badge(d["type"])}
          <div class="doc-info">
            <span class="doc-name">{d["name"]}</span>
            <span class="doc-sub">{d["chunks"]} chunks indexed</span>
          </div>
          <span style="color:{color};font-weight:700">{icon}</span>
        </div>"""
    return f'<div class="doc-list">{rows}</div>'


def _html_citations(citations: list) -> str:
    if not citations:
        return '<p class="dim">Source references appear here after each answer.</p>'
    hues  = ["#4f46e5", "#0891b2", "#ea580c", "#16a34a", "#dc2626", "#7c3aed"]
    cards = ""
    for c in citations:
        accent   = hues[(c["index"] - 1) % len(hues)]
        page_lbl = (f'<span class="cit-page">p.{c["page"]}</span>'
                    if str(c.get("page", "N/A")) != "N/A" else "")
        excerpt  = (c["content"][:800]
                    .replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))
        cards += f"""
        <div class="cit-card" style="border-left-color:{accent}">
          <div class="cit-head">
            <span class="cit-num" style="background:{accent}">[{c["index"]}]</span>
            <span class="cit-file">{c["file"]}</span>
            {page_lbl}
          </div>
          <details>
            <summary class="cit-toggle">▸  View source excerpt</summary>
            <pre class="cit-excerpt">{excerpt}</pre>
          </details>
        </div>"""
    return f'<div class="cit-list">{cards}</div>'


# ═══════════════════════════════════════════════════════════════
# 4 · RAG CHAIN  (LLM untouched — only prompt + parse layer)
# ═══════════════════════════════════════════════════════════════
RAG_PROMPT = PromptTemplate(
    input_variables=["context", "question"],
    template="""You are DocuMind AI, an expert research assistant.

## Instructions
- Answer STRICTLY from the retrieved context provided below.
- Cite every factual claim with inline references like [1], [2] matching source numbers.
- If the context does NOT contain sufficient information, output the token <<<WEB_NEEDED>>> on its own line, then provide your best general-knowledge answer.
- Keep answers clear, structured, and analytical.

## Retrieved Context
{context}

## Question
{question}

## Answer""",
)


def chat(question: str, history: list, cit_state: list):
    global vectorstore

    if not question.strip():
        return history, "", cit_state, _html_citations([])

    # ── No docs yet → web only ──────────────────────────────────────────────
    if vectorstore is None:
        try:
            web = web_search.run(question)[:900]
        except Exception:
            web = "_(web search unavailable)_"
        ans = ("⚠️ **No documents uploaded.**  "
               "Here is what I found on the web:\n\n" + web)
        history.append({"role": "user",      "content": question})
        history.append({"role": "assistant", "content": ans})
        return history, "", [], _html_citations([])

    # ── RAG chain (identical to original) ───────────────────────────────────
    retriever = vectorstore.as_retriever(search_kwargs={"k": 4})
    chain = RetrievalQA.from_chain_type(
        llm           = llm,
        retriever     = retriever,
        chain_type    = "stuff",
        chain_type_kwargs = {
            "prompt":                 RAG_PROMPT,
            "document_variable_name": "context",
        },
        return_source_documents = True,
    )
    result      = chain.invoke({"query": question})
    raw_answer  = result["result"]
    source_docs = result["source_documents"]

    # ── Web fallback ─────────────────────────────────────────────────────────
    if "<<<WEB_NEEDED>>>" in raw_answer:
        clean_ans = raw_answer.replace("<<<WEB_NEEDED>>>", "").strip()
        try:
            web_snippet = web_search.run(question)[:600]
        except Exception:
            web_snippet = "_(web search unavailable)_"
        ans = (
            "⚠️ **This answer was not found in your uploaded documents.**\n\n"
            + clean_ans
            + f"\n\n---\n🌐 **Web supplement:**\n> {web_snippet}"
        )
        citations = []

    # ── Sourced answer + citations ───────────────────────────────────────────
    else:
        ans = raw_answer
        seen, citations = set(), []
        for doc in source_docs:
            key = (doc.metadata.get("file_name"), doc.metadata.get("page"))
            if key not in seen:
                seen.add(key)
                citations.append({
                    "index":   len(citations) + 1,
                    "file":    doc.metadata.get("file_name", "Unknown"),
                    "page":    doc.metadata.get("page", "N/A"),
                    "content": doc.page_content,
                })

    history.append({"role": "user",      "content": question})
    history.append({"role": "assistant", "content": ans})
    return history, "", citations, _html_citations(citations)


# ═══════════════════════════════════════════════════════════════
# 5 · CSS  — Warm Ivory + Indigo professional light theme
# ═══════════════════════════════════════════════════════════════
CSS = """
@import url('https://fonts.googleapis.com/css2?family=Playfair+Display:wght@600;700&family=IBM+Plex+Mono:wght@400;500&family=Plus+Jakarta+Sans:wght@300;400;500;600&display=swap');

:root {
  --bg0:          #f7f6f2;
  --bg1:          #ffffff;
  --bg2:          #f0eef8;
  --bg3:          #e8e5f4;
  --border:       #ddd9f0;
  --border2:      #c9c4e8;
  --accent:       #4f46e5;
  --accent-dark:  #4338ca;
  --accent2:      #0891b2;
  --accent-light: #ede9fe;
  --fg0:          #1e1b3a;
  --fg1:          #4b4869;
  --fg2:          #8b87a8;
  --green:        #16a34a;
  --red:          #dc2626;
  --shadow:       rgba(79,70,229,.12);
  --shadow-sm:    rgba(79,70,229,.06);
  --radius:       12px;
}

*, *::before, *::after { box-sizing: border-box; }

body, .gradio-container {
  background: var(--bg0) !important;
  color:      var(--fg0) !important;
  font-family: 'Plus Jakarta Sans', sans-serif !important;
}
.gradio-container {
  max-width: 1440px !important;
  margin: 0 auto !important;
}

/* ── Header ──────────────────────────────────────────────── */
#app-header {
  text-align: center;
  padding: 38px 0 24px;
  position: relative;
}
#app-header::after {
  content: '';
  display: block;
  height: 2px;
  margin-top: 24px;
  background: linear-gradient(
    90deg, transparent 0%, var(--accent) 35%,
    var(--accent2) 65%, transparent 100%
  );
  opacity: .3;
}
.eyebrow {
  font-family: 'IBM Plex Mono', monospace;
  font-size: .68rem;
  letter-spacing: .2em;
  text-transform: uppercase;
  color: var(--accent);
  margin: 0 0 10px;
}
#app-header h1 {
  font-family: 'Playfair Display', serif;
  font-size: 2.7rem;
  color: var(--fg0);
  margin: 0 0 10px;
  letter-spacing: -.025em;
  line-height: 1.1;
}
#app-header h1 .hl {
  background: linear-gradient(115deg, var(--accent), var(--accent2));
  -webkit-background-clip: text;
  -webkit-text-fill-color: transparent;
  background-clip: text;
}
#app-header .sub {
  color: var(--fg2);
  font-size: .88rem;
  margin: 0;
  font-weight: 400;
  letter-spacing: .03em;
}

/* ── Section labels ──────────────────────────────────────── */
.panel-title {
  font-size: .67rem;
  font-weight: 700;
  letter-spacing: .18em;
  text-transform: uppercase;
  color: var(--fg2);
  margin: 0 0 14px;
  padding-bottom: 10px;
  border-bottom: 1.5px solid var(--border);
}

/* ── File uploader ───────────────────────────────────────── */
.gradio-file, label.block {
  background: var(--bg1) !important;
  border: 2px dashed var(--border2) !important;
  border-radius: var(--radius) !important;
  color: var(--fg1) !important;
  transition: border-color .2s, background .2s !important;
}
.gradio-file:hover {
  border-color: var(--accent) !important;
  background: var(--accent-light) !important;
}

/* ── Buttons — Process ───────────────────────────────────── */
#proc-btn button {
  background: var(--accent) !important;
  color: #fff !important;
  border: none !important;
  border-radius: 9px !important;
  font-weight: 600 !important;
  font-size: .88rem !important;
  letter-spacing: .02em !important;
  width: 100% !important;
  transition: background .2s, transform .12s, box-shadow .2s !important;
  box-shadow: 0 3px 14px var(--shadow) !important;
}
#proc-btn button:hover {
  background: var(--accent-dark) !important;
  transform: translateY(-1px) !important;
  box-shadow: 0 5px 18px var(--shadow) !important;
}

/* ── Status textbox ──────────────────────────────────────── */
#proc-status textarea, #proc-status input {
  background: #f0fdf4 !important;
  color: #15803d !important;
  border: 1.5px solid #bbf7d0 !important;
  border-radius: 8px !important;
  font-family: 'IBM Plex Mono', monospace !important;
  font-size: .77rem !important;
  font-weight: 500 !important;
}

/* ── Doc list ────────────────────────────────────────────── */
.doc-list { display: flex; flex-direction: column; gap: 8px; }
.doc-row  {
  display: flex;
  align-items: center;
  gap: 10px;
  background: var(--bg1);
  border: 1px solid var(--border);
  border-radius: 9px;
  padding: 10px 13px;
  transition: box-shadow .15s, border-color .15s;
}
.doc-row:hover {
  box-shadow: 0 3px 12px var(--shadow-sm);
  border-color: var(--border2);
}
.doc-info  { flex: 1; display: flex; flex-direction: column; gap: 2px; min-width: 0; }
.doc-name  {
  font-size: .83rem;
  font-weight: 600;
  color: var(--fg0);
  white-space: nowrap;
  overflow: hidden;
  text-overflow: ellipsis;
}
.doc-sub   {
  font-size: .69rem;
  color: var(--fg2);
  font-family: 'IBM Plex Mono', monospace;
}
.dim {
  color: var(--fg2);
  font-size: .83rem;
  text-align: center;
  padding: 18px 0;
  font-style: italic;
}

/* ── Citations ───────────────────────────────────────────── */
.cit-list { display: flex; flex-direction: column; gap: 10px; }
.cit-card {
  background: var(--bg1);
  border: 1px solid var(--border);
  border-left: 4px solid var(--accent);
  border-radius: 9px;
  padding: 12px 14px;
  transition: box-shadow .15s;
}
.cit-card:hover { box-shadow: 0 3px 12px var(--shadow-sm); }
.cit-head  {
  display: flex;
  align-items: center;
  gap: 8px;
  margin-bottom: 7px;
  flex-wrap: wrap;
}
.cit-num {
  background: var(--accent);
  color: #fff;
  font-family: 'IBM Plex Mono', monospace;
  font-size: .66rem;
  font-weight: 700;
  padding: 2px 9px;
  border-radius: 20px;
  white-space: nowrap;
}
.cit-file {
  font-size: .8rem;
  font-weight: 600;
  color: var(--fg1);
  flex: 1;
  word-break: break-all;
  min-width: 0;
}
.cit-page {
  font-family: 'IBM Plex Mono', monospace;
  font-size: .67rem;
  color: var(--fg2);
  background: var(--bg2);
  padding: 2px 8px;
  border-radius: 4px;
  border: 1px solid var(--border);
}
.cit-toggle {
  display: block;
  cursor: pointer;
  font-size: .77rem;
  color: var(--accent);
  font-weight: 600;
  user-select: none;
  list-style: none;
  padding: 3px 0;
  transition: color .15s;
}
.cit-toggle::-webkit-details-marker { display: none; }
details[open] .cit-toggle { color: var(--accent2); }
.cit-excerpt {
  margin: 8px 0 0;
  font-family: 'IBM Plex Mono', monospace;
  font-size: .75rem;
  color: var(--fg1);
  line-height: 1.72;
  background: var(--bg0);
  border: 1px solid var(--border);
  border-radius: 7px;
  padding: 10px 13px;
  white-space: pre-wrap;
  word-break: break-word;
  max-height: 220px;
  overflow-y: auto;
}

/* ── Chatbot shell ───────────────────────────────────────── */
.chatbot-wrap > div,
.chatbot-wrap .wrap {
  background: #faf9f6 !important;
  border: 1.5px solid var(--border) !important;
  border-radius: var(--radius) !important;
}

/* ── USER bubble ─────────────────────────────────────────── */
/* Gradio 4.x uses data-testid attributes on wrappers */
.message-bubble-border.user,
.message.user {
  background: var(--accent) !important;
  color: #ffffff !important;
  border-radius: 18px 18px 4px 18px !important;
  padding: 11px 16px !important;
  max-width: 76% !important;
  margin-left: auto !important;
  border: none !important;
  box-shadow: 0 3px 14px rgba(79,70,229,.22) !important;
  font-size: .9rem !important;
  line-height: 1.62 !important;
}
.message-bubble-border.user *,
.message.user * {
  color: #ffffff !important;
}

/* ── ASSISTANT bubble ────────────────────────────────────── */
.message-bubble-border.bot,
.message.bot {
  background: var(--bg1) !important;
  color: var(--fg0) !important;
  border: 1.5px solid var(--border) !important;
  border-radius: 18px 18px 18px 4px !important;
  padding: 13px 16px !important;
  max-width: 88% !important;
  box-shadow: 0 2px 8px rgba(0,0,0,.05) !important;
  font-size: .9rem !important;
  line-height: 1.72 !important;
}
.message-bubble-border.bot *,
.message.bot * {
  color: var(--fg0) !important;
}
/* Override for specific sub-elements */
.message.bot strong { color: var(--fg0) !important; font-weight: 700; }
.message.bot em     { color: var(--fg1) !important; }
.message.bot a      { color: var(--accent) !important; }
.message.bot code {
  background: var(--bg2) !important;
  color: var(--accent) !important;
  border: 1px solid var(--border) !important;
  border-radius: 5px;
  padding: 1px 6px;
  font-family: 'IBM Plex Mono', monospace;
  font-size: .83em;
}
.message.bot blockquote {
  border-left: 3px solid var(--accent2) !important;
  margin: 8px 0 !important;
  padding-left: 12px !important;
  color: var(--fg1) !important;
  font-style: italic;
}
.message.bot ul, .message.bot ol {
  padding-left: 20px !important;
  color: var(--fg0) !important;
}

/* ── Input textarea ──────────────────────────────────────── */
#q-input textarea {
  background: var(--bg1) !important;
  color: var(--fg0) !important;
  border: 1.5px solid var(--border2) !important;
  border-radius: 10px !important;
  font-family: 'Plus Jakarta Sans', sans-serif !important;
  font-size: .92rem !important;
  resize: none !important;
  transition: border-color .2s, box-shadow .2s !important;
}
#q-input textarea:focus {
  border-color: var(--accent) !important;
  box-shadow: 0 0 0 3px rgba(79,70,229,.1) !important;
  outline: none !important;
}
#q-input textarea::placeholder { color: var(--fg2) !important; }

/* ── Send button ─────────────────────────────────────────── */
#send-btn button {
  background: var(--accent) !important;
  color: #fff !important;
  border: none !important;
  border-radius: 10px !important;
  font-weight: 700 !important;
  font-size: .9rem !important;
  height: 100% !important;
  min-height: 58px !important;
  letter-spacing: .02em !important;
  transition: background .2s, transform .12s, box-shadow .2s !important;
  box-shadow: 0 3px 14px var(--shadow) !important;
}
#send-btn button:hover {
  background: var(--accent-dark) !important;
  transform: translateY(-1px) !important;
  box-shadow: 0 5px 18px var(--shadow) !important;
}

/* ── Action row buttons ──────────────────────────────────── */
#clear-btn button {
  background: var(--bg1) !important;
  color: var(--fg1) !important;
  border: 1.5px solid var(--border2) !important;
  border-radius: 9px !important;
  font-size: .82rem !important;
  font-weight: 500 !important;
  transition: all .2s !important;
}
#clear-btn button:hover {
  color: var(--red) !important;
  border-color: var(--red) !important;
  background: #fff5f5 !important;
}

/* ── Divider ─────────────────────────────────────────────── */
.section-divider {
  border: none;
  border-top: 1.5px solid var(--border);
  margin: 18px 0;
}

/* ── Scrollbar ───────────────────────────────────────────── */
::-webkit-scrollbar        { width: 5px; height: 5px; }
::-webkit-scrollbar-track  { background: var(--bg0); }
::-webkit-scrollbar-thumb  { background: var(--border2); border-radius: 3px; }
::-webkit-scrollbar-thumb:hover { background: var(--accent); }

/* ── Gradio generic overrides ────────────────────────────── */
.gradio-container .block { background: transparent !important; }
label > span { color: var(--fg1) !important; font-size: .8rem !important; }
"""

# ═══════════════════════════════════════════════════════════════
# 6 · GRADIO UI
# ═══════════════════════════════════════════════════════════════
with gr.Blocks(css=CSS, title="DocuMind AI", theme=gr.themes.Base()) as app:
    cit_state = gr.State([])

    # ── Header ──────────────────────────────────────────────────
    gr.HTML("""
    <div id="app-header">
      <p class="eyebrow">Intelligent Research Assistant</p>
      <h1>Docu<span class="hl">Mind</span> AI</h1>
      <p class="sub">
        Upload documents &middot; Ask natural-language questions
        &middot; Get cited, verifiable answers
      </p>
    </div>""")

    with gr.Row(equal_height=False):

        # ══════════════════════════════════
        #  LEFT SIDEBAR — Upload + Citations
        # ══════════════════════════════════
        with gr.Column(scale=1, min_width=310):

            gr.HTML('<div class="panel-title">📁  Your Documents</div>')

            file_upload = gr.File(
                label      = "Drop files here or click to browse",
                file_count = "multiple",
                file_types = [".pdf", ".docx", ".doc",
                               ".pptx", ".ppt", ".md", ".txt"],
            )
            proc_btn = gr.Button("⚙  Process & Index Files",
                                 elem_id="proc-btn", variant="primary")
            proc_status = gr.Textbox(
                label       = "",
                interactive = False,
                show_label  = False,
                elem_id     = "proc-status",
                lines       = 1,
                placeholder = "Processing status will appear here…",
            )
            doc_html = gr.HTML('<p class="dim">No documents loaded yet.</p>')

            gr.HTML('<hr class="section-divider">')

            gr.HTML('<div class="panel-title">📚  Source References</div>')
            cit_html = gr.HTML(
                '<p class="dim">Source references appear here after each answer.</p>'
            )

        # ══════════════════════════════════
        #  MAIN AREA — Chat
        # ══════════════════════════════════
        with gr.Column(scale=2):

            gr.HTML('<div class="panel-title">💬  Conversation</div>')

            chatbot = gr.Chatbot(
                label            = "",
                type             = "messages",
                height           = 520,
                show_copy_button = True,
                render_markdown  = True,
                elem_classes     = ["chatbot-wrap"],
                placeholder      = (
                    "<div style='text-align:center;padding:52px 0;color:#8b87a8'>"
                    "<div style='font-size:3rem;margin-bottom:14px'>✦</div>"
                    "<p style='font-size:1.05rem;font-weight:600;"
                    "color:#4b4869;margin:0'>"
                    "Start by uploading your documents</p>"
                    "<p style='font-size:.84rem;margin:8px 0 0'>"
                    "Supports PDF · DOCX · PPTX · Markdown · TXT</p>"
                    "</div>"
                ),
            )

            with gr.Row():
                q_input = gr.Textbox(
                    placeholder = "Ask a question about your documents…",
                    show_label  = False,
                    lines       = 2,
                    scale       = 5,
                    elem_id     = "q-input",
                )
                send_btn = gr.Button(
                    "Send  ➤", scale=1,
                    elem_id="send-btn", variant="primary", min_width=95,
                )

            with gr.Row():
                clear_btn  = gr.Button("🗑  Clear chat",
                                       scale=1, elem_id="clear-btn")

    # ═══════════════════════════════════════════════════════════
    # 7 · EVENT WIRING
    # ═══════════════════════════════════════════════════════════
    proc_btn.click(
        process_files,
        inputs  = [file_upload],
        outputs = [doc_html, proc_status],
    )

    def _chat_wrap(q, hist, cits):
        return chat(q, hist, cits)

    send_btn.click(
        _chat_wrap,
        inputs  = [q_input, chatbot, cit_state],
        outputs = [chatbot, q_input, cit_state, cit_html],
    )
    q_input.submit(
        _chat_wrap,
        inputs  = [q_input, chatbot, cit_state],
        outputs = [chatbot, q_input, cit_state, cit_html],
    )

    def _clear():
        return ([], [],
                '<p class="dim">Source references appear here after each answer.</p>')

    clear_btn.click(_clear, outputs=[chatbot, cit_state, cit_html])


# ═══════════════════════════════════════════════════════════════
# 8 · LAUNCH
# ═══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    app.launch()